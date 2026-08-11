#!/usr/bin/env python3
"""Generate end-to-end implementation slides (PPTX)."""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "docs" / "slides" / "End_to_End_Implementation.pptx"
PLOTS = ROOT / "outputs" / "plots"
METRICS = ROOT / "outputs" / "metrics" / "all_metrics.json"

NOKIA_BLUE = RGBColor(0x12, 0x41, 0x91)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NOKIA_BLUE
    shape.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = WHITE


def add_bullet_slide(prs, title, bullets: list[str], image_path: Path | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(0.8))
    tx.text_frame.paragraphs[0].text = title
    tx.text_frame.paragraphs[0].font.size = Pt(24)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.color.rgb = NOKIA_BLUE

    body_w = Inches(5.5) if image_path and image_path.exists() else Inches(9)
    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), body_w, Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(14)
        p.level = 0

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(5.8), Inches(1.2), width=Inches(4))


def add_image_slide(prs, title, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.6))
    tx.text_frame.paragraphs[0].text = title
    tx.text_frame.paragraphs[0].font.size = Pt(22)
    tx.text_frame.paragraphs[0].font.bold = True
    tx.text_frame.paragraphs[0].font.color.rgb = NOKIA_BLUE
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.5), Inches(0.9), width=Inches(9))


def main():
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    metrics = {}
    if METRICS.exists():
        metrics = json.loads(METRICS.read_text(encoding="utf-8"))

    add_title_slide(prs, "Secure Multi-Agent AI Framework\nfor RAN Control Loops",
                    "End-to-End Implementation | NBUC Project 03 | Swetha Kerahalli, Nokia")

    add_bullet_slide(prs, "Problem Statement", [
        "AI-native 6G RANs embed AI in RRC control loops — mobility, handover, resource allocation",
        "Centralized AI = single point of failure; compromised AI can manipulate RRC decisions",
        "Need: secure, distributed, verifiable multi-agent coordination with consensus",
        "References: 3GPP TS 38.331, O-RAN Near-RT RIC, IEEE MobiLLM, NBUC Problem Statement",
    ])

    add_bullet_slide(prs, "Proposed Architecture", [
        "Multi-Agent Layer: 8 specialized agents (Mobility, Security, Resource, Energy, Trust, Beamforming, QoS, Policy)",
        "RAN Control Layer: RRC decisions per TS 38.331 (HO, reconfiguration, beam switch)",
        "Security Layer: Trust engine, anomaly detection, BFT consensus, PQC (Kyber/Dilithium)",
        "O-RAN Integration: Near-RT RIC, E2 interface, xApps/rApps",
        "Digital Twin: 48 cells, 12 gNBs, attack injection & mitigation simulation",
    ], PLOTS / "architecture" / "01_overall_evaluation.png")

    add_bullet_slide(prs, "Dataset & Data Split", [
        "Synthetic telemetry: 70,000 rows × 42 columns",
        "6 scenarios: normal, jamming, compromised agent, adversarial mobility, massive mobility, congestion",
        "RAG corpus: 70 knowledge chunks (3GPP, O-RAN, MARL, security)",
        "Split: 70% train / 15% validation / 15% test (stratified by scenario)",
        "Features: RSRP, RSRQ, SINR, CQI, PRB, trust, anomaly, agent confidence",
    ], PLOTS / "data" / "01_dataset_overview.png")

    add_image_slide(prs, "Data Correlation Heatmap & CDFs", PLOTS / "data" / "02_correlation_heatmap.png")

    add_bullet_slide(prs, "Agent Models — Training / Validation / Testing", [
        "Mobility Agent: Gradient Boosting → ho_required (F1 on test)",
        "Security Agent: Random Forest → threat_label (benign/malicious)",
        "Resource Agent: Gradient Boosting → allocated_prb_count (regression)",
        "Trust Agent: MLP → trust_score regression",
        "QoS/Policy/Energy/Beamforming: RF/GB classifiers & regressors",
        "All models: train/val/test metrics saved; learning curves & confusion matrices generated",
    ], PLOTS / "agents" / "03_agent_performance_heatmap.png")

    if metrics.get("agents"):
        agent_lines = []
        for name, ag in metrics["agents"].items():
            key = "f1" if ag["task"] == "classification" else "r2"
            agent_lines.append(f"{name}: test {key.upper()}={ag['test'].get(key, 0):.3f}")
        add_bullet_slide(prs, "Test Set Results (Sample)", agent_lines[:8])

    add_image_slide(prs, "Classification Agent — Confusion Matrix & ROC", PLOTS / "agents" / "mobility_agent_confusion_matrix.png")
    add_image_slide(prs, "Digital Twin — Time Series KPIs", PLOTS / "digital_twin" / "01_twin_time_series.png")
    add_image_slide(prs, "Digital Twin — Cell State Map", PLOTS / "digital_twin" / "02_twin_cell_map.png")

    add_bullet_slide(prs, "Consensus & Orchestration", [
        "Weighted trust voting with BFT-inspired validation",
        "Thresholds: majority >70%, trust >0.8, confidence >0.85",
        "MultiAgentCoordinator runs batch inference on test set",
        f"Consensus accept rate: {metrics.get('coordinator', {}).get('consensus_accept_rate', 0):.1%}",
        "Rejected proposals from compromised/low-trust agents",
    ])

    add_bullet_slide(prs, "Dashboard & Chatbot", [
        "Streamlit dashboard: Overview, Data, Agents, Train/Val/Test, Digital Twin, Chatbot",
        "Plotly interactive charts + static matplotlib plots",
        "RAG chatbot over project knowledge corpus (TF-IDF retrieval)",
        "Launch: streamlit run dashboard/app.py",
    ])

    add_bullet_slide(prs, "References", [
        "3GPP TS 38.300, 38.331, 38.215, 28.530, 23.288, 33.501",
        "3GPP TR 38.817, TR 23.700-80",
        "O-RAN Near-RT RIC, E2GAP, xApp/rApp specifications",
        "IEEE: MobiLLM, AI-Augmented Predictive Mobility, Jamming-Resilient HO (RL)",
        "NIST PQC: CRYSTALS-Kyber, Dilithium | NBUC Problem Statement — Swetha Kerahalli",
    ])

    add_title_slide(prs, "Thank You", "Questions? | proj3_code/ | Dashboard + Docs + Models")

    prs.save(str(OUTPUT))
    print(f"Slides saved: {OUTPUT}")


if __name__ == "__main__":
    main()
